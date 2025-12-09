from pathlib import Path
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import io

def extract_text_from_file(uploaded_file) -> str:
    """Return plain text from an uploaded Streamlit file."""
    suffix = Path(uploaded_file.name).suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(uploaded_file)
    elif suffix == ".pptx":
        return _extract_pptx(uploaded_file)
    elif suffix == ".docx":
        return _extract_docx(uploaded_file)
    elif suffix in [".txt"]:
        return uploaded_file.read().decode("utf-8", errors="ignore")
    else:
        return "Unsupported file type" 


def _extract_pdf(uploaded_file) -> str:
    reader = PdfReader(uploaded_file)
    text = []
    for page in reader.pages:
        text.append(page.extract_text() or "")
    return "\n".join(text)

def _extract_pptx(uploaded_file) -> str:
    presentation = Presentation(uploaded_file)
    text_runs = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def _extract_docx(uploaded_file) -> str:
    file_bytes = uploaded_file.read()
    doc = docx.Document(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs)
